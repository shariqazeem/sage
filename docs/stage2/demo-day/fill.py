"""Fill the official OpenClaw Demo Day template with Sage's content.

Rule: never assign text_frame.text (that collapses the paragraph to one unstyled run and
loses the template's typography). Set the FIRST run's text and blank the rest, so every
size, weight and colour the designers chose survives.
"""
from pptx import Presentation
import copy, sys

prs = Presentation("template.pptx")
slides = list(prs.slides)
missed = []

def find(slide, needle):
    for sh in slide.shapes:
        if sh.has_text_frame and needle.strip().lower() in sh.text_frame.text.strip().lower():
            return sh
    return None

def put(slide, needle, new):
    """Replace the text of the box containing `needle`. `new` may be a list of paragraphs."""
    sh = find(slide, needle)
    if sh is None:
        missed.append(needle[:48]); return None
    tf = sh.text_frame
    paras = new if isinstance(new, list) else [new]
    # reuse existing paragraphs, cloning the last one's formatting when we need more
    while len(tf.paragraphs) < len(paras):
        tf._txBody.append(copy.deepcopy(tf.paragraphs[-1]._p))
    for i, p in enumerate(tf.paragraphs):
        if i < len(paras):
            if not p.runs:
                continue
            p.runs[0].text = paras[i]
            for r in p.runs[1:]:
                r.text = ""
        else:
            for r in p.runs:
                r.text = ""
    return sh

def drop(slide, needle):
    sh = find(slide, needle)
    if sh is None:
        missed.append("DROP " + needle[:44]); return
    sh._element.getparent().remove(sh._element)

# ─────────────────────────────────────────── 1 · TITLE
s = slides[0]
put(s, "[YOUR PROJECT / AGENT NAME HERE]", "Sage")
put(s, "Insert One-Sentence Pitch",
    '"An AI agent that hires real people to test your product, and pays them in USDC."')
put(s, "[Your Name / Team Handles]", "Shariq  ·  @sagepaysai")
put(s, "[AI Agent / x402 / DeFI / Infra]", "AI Agent  ·  x402  ·  ERC-8004")

# ─────────────────────────────────────────── 2 · TABLE OF CONTENTS
s = slides[1]
put(s, "Hook the audience, state your Agent",
    "What Sage is, and the five steps it runs with no human in any of them.")
put(s, "Explain how your product/agent solves the problem using Metis",
    "Founders ship products nobody has used. Every existing option makes them do the work themselves.")
put(s, "The heart of your pitch! Show your agent live",
    "A product it has never seen: it browses, writes missions, pays a real tester — and refuses one.")
put(s, "Highlight real users, transaction metrics",
    "$48.10 settled to 16 people on GOAT mainnet. 48% of submissions refused. Every payout public.")
put(s, "Golden Pitch Rule",
    "Everything in this deck resolves to a public transaction on GOAT mainnet. You do not need our cooperation to check a single number in it.")

# ─────────────────────────────────────────── 3 · PROJECT & AGENT OVERVIEW
s = slides[2]
put(s, "[Agent Name]", "Sage")
put(s, "Insert Catchy Tagline",
    '"Hire an AI worker. Give it a budget, not your keys."')
put(s, "Builder Guidance Checklist", "What the agent actually does")
put(s, "Keep your intro under 45 seconds",
    "1.  Opens your product in a real browser and explores it.")
put(s, "State clearly whether your project is an autonomous agent",
    "2.  Writes the testing missions from what it observed.   3.  Real people complete them.")
put(s, "Avoid deep technical jargon upfront",
    ["4.  Checks every report against its own browsing.",
     "5.  Pays in USDC — or refuses, with a reason.",
     "",
     "The founder only approves the plan and funds it."])
put(s, "Tip for Builders",
    "Sage is an autonomous agent, not a tool a human drives. It decides who gets paid, and a smart contract — not a prompt — decides how much.")

# ─────────────────────────────────────────── 4 · PROBLEM & MARKET NEED
s = slides[3]
put(s, "The Current Friction", "The Current Friction")
put(s, "Describe the pain point users face today",
    "You shipped it. Nobody has used it. An agency costs $5,000. A bug bounty returns noise. Friends are only ever kind.")
put(s, "How [Agent Name] Solves It", "How Sage Solves It")
put(s, "Explain how your product or agent automates",
    "A URL and a budget. It explores, writes the missions, judges the testers, and pays from a vault it cannot exceed.")
put(s, "Guidance: Focus on 1 strong problem-solution contrast",
    "One contrast: today the founder recruits, reviews and pays. With Sage they approve a plan and fund a vault.")

# ─────────────────────────────────────────── 5 · ARCHITECTURE & CORE FEATURES
s = slides[4]
put(s, "1. Autonomous Execution", "1. Autonomous Execution")
put(s, "Describe core AI logic, prompt workflows",
    "Three model layers: one designs missions, one judges evidence, one talks to founders.")
put(s, "2. Metis & OpenClaw Stack", "2. GOAT, Metis & OpenClaw Stack")
put(s, "Mention ERC-8004 identity, x402 micro-payments",
    "ERC-8004 identity #79, x402 payments, and a CampaignVault live on GOAT mainnet.")
put(s, "3. User UX & Outcome", "3. Bounded Autonomy Over Money")
put(s, "Highlight end-user benefit",
    "The vault computes every reward and enforces every cap. No model here computes an amount.")
put(s, "Technical Pitch Tip",
    "Two front doors, one engine: a browser wallet on sagepays.xyz, or the whole loop walletless from Telegram.")

# ─────────────────────────────────────────── 6 · LIVE PRODUCT DEMO
s = slides[5]
put(s, "[SWITCH TO LIVE PRODUCT / APP SCREEN]", "[SWITCH TO LIVE PRODUCT — sagepays.xyz]")
put(s, "Or embed a high-resolution GIF",
    "A product Sage has never inspected, a real payout on GOAT mainnet, and a submission it refuses.")
put(s, "Live Demo Step-By-Step Plan", "Live Demo Step-By-Step Plan")
put(s, "Step 1 (0:30s)",
    "Step 1 (0:35s): A product it has never seen. It browses, then writes the missions itself.")
put(s, "Step 2 (0:45s)",
    "Step 2 (0:30s): A real tester is paid in USDC — receipt and transaction on screen.")
put(s, "Step 3 (0:45s)",
    "Step 3 (0:40s): I try to get paid without doing the work. Sage refuses me, and says why.")

# ─────────────────────────────────────────── 7 · TRACTION
s = slides[6]
put(s, "[NUM]", "16")
put(s, "Real Active Users / Interactions", "People paid in real USDC")
put(s, "[NUM]", "$48.10")
put(s, "Volume / x402 Payments Processed", "USDC settled on GOAT mainnet")
put(s, "[PERCENT]", "48%")
put(s, "Agent Execution Success Rate", "Submissions the agent refused")
put(s, "What We Shipped in 8 Weeks:", "What we shipped:")
put(s, "Successfully published and verified AI Agent on ClawUp",
    "68 products inspected and 42 tester reports judged — 22 paid, 20 refused, none left unresolved.")
put(s, "Integrated x402 micro-payment layer and ERC-8004 identity",
    "ERC-8004 identity, x402 rail and a V2 CampaignVault live on GOAT mainnet; MCP server built for ClawUp.")
put(s, "Acquired initial organic users via GEO growth",
    "A paid campaign onboarded 6 people to ClawUp; 2nd place cohort-wide in #ClawToTheTop.")
put(s, "Traction Tip",
    "Median time from a stranger submitting evidence to USDC arriving in their wallet: 175 seconds. Fastest: 15 seconds.")

# ─────────────────────────────────────────── 8 · ROADMAP
s = slides[7]
put(s, "Shipped working AI Agent, integrated Metis & ClawUp stack.",
    "Autonomous payouts live on GOAT mainnet. 16 people paid, 20 refused.")
put(s, "[List key feature upgrades, mainnet rollout, user growth goal]",
    "Outside founders funding their own vaults, starting with grant applicants.")
put(s, "[GOAT Grant funding, deeper Metis L2 scaling",
    "Agent-to-agent testing over MCP and x402, effort-weighted rewards, and GOAT grant participation.")
put(s, "Explain briefly how your project contributes long-term value",
    "Every campaign puts real USDC through GOAT and leaves a public receipt. Supply is proven; the work now is founders with budgets.")
put(s, "Closing Pitch Tip",
    "Most agents talk. Sage decides — and pays. That is the category we intend to own.")

# ─────────────────────────────────────────── 9 · THANK YOU
s = slides[8]
put(s, "@[YourProjectHandle]", "@sagepaysai")
put(s, "[your-agent.clawup.ai]", "sagepays.xyz")
put(s, "github.com/[team/repo]", "github.com/shariqazeem/sage")
put(s, "Try [Agent Name] live on ClawUp today!", "Every payout has a public receipt. Go and check one.")

prs.save("Sage-DemoDay-Official.pptx")
print("saved Sage-DemoDay-Official.pptx")
if missed:
    print("\n!! NOT FOUND (check these):")
    for m in missed: print("   -", m)
