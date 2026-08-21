"""Fix what only showed up once the deck was actually rendered and looked at.

The template bakes placeholder copy INTO images on slides 3 and 4, so text replacement
could never reach it — those strips are rebuilt as real shapes in the template's own tints
(sampled from the images themselves). The rest is dark-on-dark text, a row of ticks with
nothing beside one of them, and two boxes running off their card.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
E = 914400.0
prs = Presentation("Sage-DemoDay-Official.pptx")
S = list(prs.slides)
INK   = RGBColor(0x1A,0x1D,0x21); MID = RGBColor(0x4A,0x4F,0x57)
TERRA = RGBColor(0xC2,0x41,0x0C); WHITE = RGBColor(0xFF,0xFF,0xFF)

def pic_at(slide, x, y, tol=0.04):
    for sh in slide.shapes:
        if sh.__class__.__name__ == "Picture" and abs(sh.left/E-x) < tol and abs(sh.top/E-y) < tol:
            return sh
    return None

def tb_with(slide, needle):
    for sh in slide.shapes:
        if sh.has_text_frame and needle.lower() in sh.text_frame.text.lower():
            return sh
    return None

def strip(slide, x, y, w, h, fill, label, body, size=11):
    """Rebuild one of the template's baked-in placeholder strips as real, editable content."""
    old = pic_at(slide, x, y)
    if old is not None:
        old._element.getparent().remove(old._element)
    r = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    r.adjustments[0] = 0.12
    r.fill.solid(); r.fill.fore_color.rgb = fill
    r.line.fill.background(); r.shadow.inherit = False
    tf = r.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.14)
    tf.margin_top = tf.margin_bottom = Inches(0.04)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    a = p.add_run(); a.text = label
    a.font.bold = True; a.font.size = Pt(size); a.font.color.rgb = INK; a.font.name = "Calibri"
    b = p.add_run(); b.text = body
    b.font.size = Pt(size); b.font.color.rgb = MID; b.font.name = "Calibri"
    return r

def settext(sh, s, size=None, color=None, bold=None):
    tf = sh.text_frame
    for i, p in enumerate(tf.paragraphs):
        for j, r in enumerate(p.runs):
            r.text = s if (i == 0 and j == 0) else ""
            if size:  r.font.size = Pt(size)
            if color: r.font.color.rgb = color
            if bold is not None: r.font.bold = bold

# ── SLIDE 3 ────────────────────────────────────────────────────────────────
s = S[2]
strip(s, 0.84, 2.66, 5.40, 0.69, RGBColor(0xF9,0xFA,0xFB),
      "Target User:  ",
      "Early founders with a live product and not enough real users to know what breaks.")
strip(s, 0.84, 3.47, 5.40, 0.69, RGBColor(0xF9,0xFA,0xFB),
      "Core Stack:  ",
      "GOAT Network mainnet · ERC-8004 identity #79 · x402 · CampaignVault")

# the placeholder badge on the card
b = pic_at(s, 4.77, 1.55)
if b is not None:
    b._element.getparent().remove(b._element)
badge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.77), Inches(1.55), Inches(1.56), Inches(0.26))
badge.adjustments[0] = 0.5
badge.fill.solid(); badge.fill.fore_color.rgb = TERRA
badge.line.fill.background(); badge.shadow.inherit = False
tf = badge.text_frame; tf.word_wrap = False
tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "LIVE ON MAINNET"
r.font.size = Pt(9); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Calibri"

# four tick rows, four lines — one of them was standing next to nothing
rows = [
    ("Opens your product in a real browser and explores it.",        2.04, 0.36),
    ("Writes the testing missions from what it observed.",           2.40, 0.55),
    ("Real people complete them; it checks every report against its own browsing.", 2.99, 0.55),
    ("Pays in USDC — or refuses, with a written reason.",            3.58, 0.55),
]
existing = [tb_with(s, "1.  Opens your product"), tb_with(s, "2.  Writes the testing"), tb_with(s, "4.  Checks every report")]
for sh, (txt, y, h) in zip(existing, rows[:3]):
    if sh is None: continue
    settext(sh, txt, size=13, color=INK, bold=False)
    sh.top, sh.height, sh.left, sh.width = Inches(y), Inches(h), Inches(7.32), Inches(5.18)
# the fourth row needs a box of its own
import copy
src = existing[0]
new = copy.deepcopy(src._element); src._element.addnext(new)
four = [sh for sh in s.shapes if sh.has_text_frame and sh.text_frame.text.strip() == rows[0][0]][-1]
settext(four, rows[3][0], size=13, color=INK, bold=False)
four.top, four.height = Inches(rows[3][1]), Inches(rows[3][2])
# and the closing line, clear of the ticks
tail = s.shapes.add_textbox(Inches(7.32), Inches(4.30), Inches(5.18), Inches(0.6))
tf = tail.text_frame; tf.word_wrap = True
tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
p = tf.paragraphs[0]
r = p.add_run(); r.text = "The founder only approves the plan and funds it."
r.font.size = Pt(13); r.font.italic = True; r.font.color.rgb = MID; r.font.name = "Calibri"

# ── SLIDE 4 ────────────────────────────────────────────────────────────────
s = S[3]
strip(s, 0.84, 3.21, 5.40, 0.58, RGBColor(0xFE,0xF2,0xF2),
      "Market Gap:  ",
      "Founders pay for testing they cannot verify, or get feedback they cannot trust.", size=10.5)
strip(s, 7.09, 3.21, 5.40, 0.58, RGBColor(0xF0,0xFD,0xF4),
      "Sage's Edge:  ",
      "Every payout is judged against the agent's own browsing and settles on-chain.", size=10.5)

# ── SLIDE 6 ────────────────────────────────────────────────────────────────
s = S[5]
for needle, size in (("[SWITCH TO LIVE PRODUCT", 14), ("A product Sage has never inspected", 11.5)):
    sh = tb_with(s, needle)
    if sh is not None:
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                r.font.color.rgb = WHITE      # it was dark ink on a dark mockup
                r.font.size = Pt(size)
# the step text was sitting on top of its own number
for needle in ("Step 1 (0:35s)", "Step 2 (0:30s)", "Step 3 (0:40s)"):
    sh = tb_with(s, needle)
    if sh is not None:
        sh.left, sh.width = Inches(7.46), Inches(5.04)
for lbl, y in (("1.", 2.06), ("2.", 2.72), ("3.", 3.38)):
    for sh in s.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip() == lbl and sh.width < Inches(0.4):
            sh.left, sh.width, sh.top, sh.height = Inches(7.15), Inches(0.28), Inches(y), Inches(0.30)
# the guidance bar on this slide was left empty
tip = s.shapes.add_textbox(Inches(1.14), Inches(6.33), Inches(9.5), Inches(0.32))
tf = tip.text_frame; tf.word_wrap = True
tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
p = tf.paragraphs[0]
r = p.add_run(); r.text = "Two minutes of the five are this demo — a real payout, and a real refusal, both on GOAT mainnet."
r.font.size = Pt(12.5); r.font.bold = True; r.font.color.rgb = INK; r.font.name = "Calibri"

# ── SLIDE 9 ────────────────────────────────────────────────────────────────
s = S[8]
gh = tb_with(s, "github.com/shariqazeem")
if gh is not None:
    gh.left, gh.width = Inches(8.34), Inches(3.05)
me = tb_with(s, "Metis Ecosystem")
if me is not None:
    me.left, me.width, me.top = Inches(12.30), Inches(0.95), Inches(6.93)
    for p in me.text_frame.paragraphs:
        for r in p.runs: r.font.size = Pt(8)

prs.save("Sage-DemoDay-Official.pptx")
print("polished")
