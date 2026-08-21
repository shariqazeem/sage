"""Render the pptx to HTML at true geometry so it can actually be looked at.

No LibreOffice here, so this places every shape at its real inches-position with its real
font size. It is an approximation of PowerPoint's layout, but it is exact about geometry —
which is what overlap and alignment defects are made of.
"""
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.text import PP_ALIGN
import base64, os, html
E = 914400.0
SCALE = 96  # px per inch
prs = Presentation("Sage-DemoDay-Official.pptx")

def color_of(run):
    try:
        c = run.font.color
        if c and c.type is not None and c.rgb is not None:
            return "#" + str(c.rgb)
    except Exception:
        pass
    return "#1a1d21"

out = ["""<meta charset="utf-8"><style>
body{background:#3a3a3a;margin:0;padding:24px;font-family:Calibri,Arial,sans-serif}
.slide{position:relative;width:%dpx;height:%dpx;background:#fff;margin:0 auto 28px;
 overflow:hidden;box-shadow:0 4px 18px rgba(0,0,0,.5)}
.n{position:absolute;left:0;top:-22px;color:#eee;font:600 14px Arial}
.tb{position:absolute;box-sizing:border-box;line-height:1.22}
img{position:absolute}
</style>""" % (13.333*SCALE, 7.5*SCALE)]

for i, s in enumerate(prs.slides, 1):
    out.append(f'<div style="position:relative;width:{13.333*SCALE}px;margin:0 auto 46px"><div class="n">SLIDE {i}</div><div class="slide">')
    for sh in s.shapes:
        try:
            x, y = sh.left/E*SCALE, sh.top/E*SCALE
            w, h = sh.width/E*SCALE, sh.height/E*SCALE
        except Exception:
            continue
        if sh.shape_type == 13 or sh.__class__.__name__ == "Picture":
            try:
                b64 = base64.b64encode(sh.image.blob).decode()
                out.append(f'<img src="data:image/png;base64,{b64}" style="left:{x}px;top:{y}px;width:{w}px;height:{h}px">')
            except Exception:
                pass
        elif sh.has_text_frame and sh.text_frame.text.strip():
            paras = []
            for p in sh.text_frame.paragraphs:
                al = {PP_ALIGN.CENTER: "center", PP_ALIGN.RIGHT: "right"}.get(p.alignment, "left")
                runs = []
                for r in p.runs:
                    if not r.text: continue
                    sz = r.font.size.pt if r.font.size else 12
                    wt = "700" if r.font.bold else "400"
                    it = "italic" if r.font.italic else "normal"
                    runs.append(f'<span style="font-size:{sz*SCALE/72:.1f}px;font-weight:{wt};'
                                f'font-style:{it};color:{color_of(r)}">{html.escape(r.text)}</span>')
                paras.append(f'<div style="text-align:{al}">{"".join(runs) or "&nbsp;"}</div>')
            out.append(f'<div class="tb" style="left:{x}px;top:{y}px;width:{w}px;height:{h}px;'
                       f'outline:1px dashed rgba(255,0,0,.35)">{"".join(paras)}</div>')
    out.append("</div></div>")
open("preview.html","w").write("\n".join(out))
print("wrote preview.html")
