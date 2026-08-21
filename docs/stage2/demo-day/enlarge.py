"""Grow the body text so it survives being watched on a laptop.

The template sets body copy at 11.25pt, which is fine on a printed page and far too small
in a 5-minute video. Its white cards are much larger than the text boxes sitting on them,
so the room is already there — this takes it.
"""
from pptx import Presentation
from pptx.util import Pt, Inches
E = 914400.0
prs = Presentation("Sage-DemoDay-Official.pptx")
S = list(prs.slides)

def box(slide, needle):
    for sh in slide.shapes:
        if sh.has_text_frame and needle.strip().lower() in sh.text_frame.text.strip().lower():
            return sh
    return None

def grow(slide, needle, pt=None, h=None, top=None, left=None):
    sh = box(slide, needle)
    if sh is None:
        print("  !! not found:", needle[:44]); return
    if pt:
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(pt)
    if h:    sh.height = Inches(h)
    if top:  sh.top = Inches(top)
    if left: sh.left = Inches(left)
    sh.text_frame.word_wrap = True

# S2 · contents
for n in ["What Sage is, and the five", "Founders ship products nobody",
          "A product it has never seen:", "$48.10 settled to 16 people"]:
    grow(S[1], n, pt=13, h=1.05)
grow(S[1], "Everything in this deck resolves", pt=13, h=0.5, top=6.13)

# S3 · the five steps  (card runs to y=5.33, so there is room)
grow(S[2], "1.  Opens your product", pt=14.5, h=0.42, top=2.06)
grow(S[2], "2.  Writes the testing missions", pt=14.5, h=0.70, top=2.56)
grow(S[2], "4.  Checks every report", pt=14.5, h=1.60, top=3.34)
grow(S[2], "Sage is an autonomous agent", pt=12.5, h=0.32, top=6.33)

# S4 · problem / solution  (cards run to y=4.07)
grow(S[3], "You shipped it. Nobody has used it", pt=15, h=1.15, top=2.62)
grow(S[3], "A URL and a budget. It explores", pt=15, h=1.15, top=2.62)
grow(S[3], "One contrast: today the founder", pt=12.5, h=0.32, top=6.33)

# S5 · three columns  (cards run to y=4.07)
for n in ["Three model layers:", "ERC-8004 identity #79", "The vault computes every reward"]:
    grow(S[4], n, pt=13.5, h=0.68, top=3.32)
grow(S[4], "Two front doors, one engine", pt=12.5, h=0.32, top=6.33)

# S6 · demo steps  (card runs to y=5.02) — respace so the bigger type cannot collide
grow(S[5], "Step 1 (0:35s)", pt=13.5, h=0.55, top=2.06)
grow(S[5], "Step 2 (0:30s)", pt=13.5, h=0.55, top=2.72)
grow(S[5], "Step 3 (0:40s)", pt=13.5, h=0.55, top=3.38)
for lbl, y in (("1.", 2.06), ("2.", 2.72), ("3.", 3.38)):
    sh = box(S[5], lbl)
    if sh is not None and sh.width < Inches(0.4):
        for p in sh.text_frame.paragraphs:
            for r in p.runs: r.font.size = Pt(13.5)
        sh.top = Inches(y); sh.height = Inches(0.3)
grow(S[5], "A product Sage has never inspected", pt=12, h=0.6)

# S7 · traction
for n in ["People paid in real USDC", "USDC settled on GOAT mainnet", "Submissions the agent refused"]:
    grow(S[6], n, pt=13, h=0.28)
for n in ["68 products inspected", "ERC-8004 identity, x402 rail", "A paid campaign onboarded"]:
    grow(S[6], n, pt=14)
grow(S[6], "Median time from a stranger", pt=12.5, h=0.32, top=6.33)

# S8 · roadmap  (the band below starts at y=3.25)
for n in ["Autonomous payouts live on GOAT", "Outside founders funding their own",
          "Agent-to-agent testing over MCP"]:
    grow(S[7], n, pt=12.5, h=0.70, top=2.50)
grow(S[7], "Every campaign puts real USDC", pt=13, h=0.30)
grow(S[7], "Most agents talk. Sage decides", pt=12.5, h=0.32, top=6.33)

prs.save("Sage-DemoDay-Official.pptx")
print("enlarged")
