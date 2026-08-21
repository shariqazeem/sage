from pptx import Presentation
p = Presentation("Sage-DemoDay-Official.pptx")
N = [
# 1
"""[0:00-0:25]  TITLE

Hi, I'm Shariq, and this is Sage.

Sage is an AI agent that hires real people to test your product, and pays them in USDC.

Not a prototype. Forty-eight dollars has already moved on GOAT mainnet to sixteen real people, and the agent refused twenty of the forty-two reports it was sent.""",
# 2
"""[0:25-0:40]  CONTENTS

Four things in five minutes: what Sage is, the problem it solves, a live demo on a product it has never seen, and where we actually got to.

One thing to say up front — every number in this deck resolves to a public transaction on GOAT mainnet. You don't need my cooperation to check any of it.""",
# 3
"""[0:40-1:15]  OVERVIEW

Sage is an autonomous agent, not a tool a human drives.

Five steps. It opens your product in a real browser and explores it. It writes the testing missions from what it observed - not from your marketing copy. Real people complete them. It checks each report against its own browsing. Then it pays in USDC, or it refuses.

The founder does two things in that entire loop: approve the plan, and fund it. A human is in none of the five steps.""",
# 4
"""[1:15-1:45]  PROBLEM

Building stopped being the hard part. Getting real strangers to use your product and tell you the truth is now the hard part.

An agency costs five thousand dollars and two weeks. A bug bounty gives you noise and duplicates and says nothing about whether onboarding makes sense. And posting "please try my app" gets you friends being kind to you.

All three make you do the work - find the people, brief them, check their work, pay them one by one.

Sage takes a URL and a budget, and does all four itself.""",
# 5
"""[1:45-2:05]  ARCHITECTURE

Three separate model layers - one designs missions, one judges evidence, one talks to founders. They're split on purpose, so a conversation can never reach the layer that decides payouts.

We're on ERC-8004 identity number seventy-nine, x402 for payments, and a CampaignVault live on GOAT mainnet.

And the important part: the vault computes every reward and enforces every cap. No model in this system computes a money amount. If one got jailbroken tomorrow, it still couldn't move a dollar.

The agent proposes. The vault disposes.""",
# 6
"""[2:05-4:05]  LIVE DEMO  -  2 minutes, 40% of the pitch

[SWITCH TO SCREEN RECORDING]

Step 1 (0:35) - "I've never pointed it at this before." Paste the URL. Let the browsing play with NO narration. Then scroll the plan: "I didn't write any of this. It wrote these missions from what it just saw." Point at the budget split.

Step 2 (0:30) - A real tester submits. Stay quiet while it judges. When it settles: "Paid. Real USDC on GOAT mainnet. That's the transaction hash, and this receipt shows the exact evidence it paid for."

Step 3 (0:40) - "Now let me try to cheat it." Submit something thin from a second wallet. When it refuses, read the reason out loud: "Refused. And it tells me why. I built this thing and it just said no to me."

Then cut to the GOAT explorer: "That's the block explorer. Not my dashboard - the chain."

[15s optional] Telegram: "Same agent, same money, from a chat. No wallet app anywhere." """,
# 7
"""[4:05-4:30]  TRACTION

Sixteen people paid in real USDC. Forty-eight dollars ten settled on GOAT mainnet.

And the number I care about most: forty-eight percent of submissions were refused.

Anyone can build an agent that pays. Building one that looks at real work and says no is the whole difficulty - and it's what makes the payouts mean anything.

Sixty-eight products inspected, forty-two reports judged, nothing left unresolved. Median time from a stranger submitting to USDC landing in their wallet: a hundred and seventy-five seconds.""",
# 8
"""[4:30-4:50]  ROADMAP

Supply is solved. Ten strangers did paid testing work in forty-five minutes. The constraint was never finding testers - it's founders with budgets.

So the next thirty days is outside founders funding their own vaults, starting with builders who need real user evidence for grant applications.

After that: agent-to-agent. Sage is already callable over MCP and x402, so other agents can commission testing with no human at all.""",
# 9
"""[4:50-5:00]  CLOSE

Most agents talk. Sage decides - and pays.

It's live at sagepays.xyz, and every payout has a public receipt. Go and check one.

Thank you.""",
]
for s, t in zip(p.slides, N):
    s.notes_slide.notes_text_frame.text = t
p.save("Sage-DemoDay-Official.pptx")
print("notes added to", len(N), "slides")
